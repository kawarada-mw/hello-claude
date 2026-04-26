#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
mamawell 社員証カード作成スクリプト
- A4横向き、6枚レイアウト (3列 x 2行)
- 角丸断裁前提のため、デザインは長方形で出力
- ロゴ画像 logo.png があれば自動埋め込み

使用ファイル（任意）:
  /home/user/hello-claude/logo.png   ← mamawellロゴ画像
"""

import os
from pptx import Presentation
from pptx.util import Mm, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ===== 色 =====
RED       = RGBColor(0xFD, 0x70, 0x66)   # #fd7066
PINK      = RGBColor(0xFB, 0xDF, 0xD2)   # #FBDFD2
BLUE_GRAY = RGBColor(0xDC, 0xE7, 0xE9)   # #DCE7E9
DARK      = RGBColor(0x4E, 0x4B, 0x4C)   # #4E4B4C
LGRAY     = RGBColor(0xD6, 0xD6, 0xD6)   # #D6D6D6
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)

# ===== 寸法 =====
CARD_W, CARD_H = Mm(85.6), Mm(54)
PAGE_W, PAGE_H = Mm(297), Mm(210)
ROWS, COLS     = 2, 3
GAP_X, GAP_Y   = Mm(5), Mm(5)
MARGIN_X = (PAGE_W - CARD_W * COLS - GAP_X * (COLS - 1)) / 2
MARGIN_Y = (PAGE_H - CARD_H * ROWS - GAP_Y * (ROWS - 1)) / 2

LOGO_PATH = '/home/user/hello-claude/logo.png'


# ===== 共通ユーティリティ =====
def no_line(shape):
    shape.line.fill.background()


def add_filled_oval(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    no_line(s)
    return s


def add_text(slide, x, y, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT,
             font_name=None, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = size
    r.font.bold = bold
    r.font.color.rgb = color
    if font_name:
        r.font.name = font_name
    return box


# ===== カードのパーツ =====
def card_background(slide, x, y):
    """白い角丸無しのカード背景（断裁前提なので長方形）"""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, CARD_W, CARD_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.color.rgb = LGRAY
    bg.line.width = Pt(0.25)


def decorations(slide, x, y):
    """背景の装飾シェイプ群（カード内に収める）"""
    # 右上: ピンクの雲型（重なる楕円）
    add_filled_oval(slide, x + Mm(60), y - Mm(2),  Mm(30), Mm(18), PINK)
    add_filled_oval(slide, x + Mm(72), y + Mm(10), Mm(20), Mm(13), PINK)
    add_filled_oval(slide, x + Mm(78), y + Mm(0),  Mm(12), Mm(8),  PINK)

    # 左下: 薄ブルーグレーの雲型
    add_filled_oval(slide, x - Mm(2),  y + Mm(42), Mm(30), Mm(15), BLUE_GRAY)
    add_filled_oval(slide, x + Mm(15), y + Mm(46), Mm(20), Mm(10), PINK)
    add_filled_oval(slide, x + Mm(2),  y + Mm(48), Mm(15), Mm(8),  BLUE_GRAY)

    # 下部: 赤い波線（滑らかな曲線）
    draw_wave(slide, x, y)


def draw_wave(slide, x, y):
    """下部の赤い細い波線"""
    sp = slide.shapes.build_freeform(x, y + Mm(48), scale=1.0)
    pts = [
        (Mm(10),  Mm(-2)),
        (Mm(20),  Mm(2)),
        (Mm(30),  Mm(0)),
        (Mm(42),  Mm(-3)),
        (Mm(54),  Mm(-2)),
        (Mm(65),  Mm(2)),
        (Mm(75),  Mm(1)),
        (Mm(85.6), Mm(-1)),
    ]
    sp.add_line_segments(pts, close=False)
    line = sp.convert_to_shape()
    line.fill.background()
    line.line.color.rgb = RED
    line.line.width = Pt(1.0)


def logo_area(slide, x, y):
    """左側のロゴエリア。ロゴ画像があれば挿入、なければシェイプで近似"""
    area_left = x + Mm(2)
    area_top  = y + Mm(8)
    area_w    = Mm(33)
    area_h    = Mm(35)

    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(
            LOGO_PATH, area_left, area_top, width=area_w, height=area_h
        )
        return

    # ===== ロゴ画像が無い場合の代替表現 =====
    cx = x + Mm(17)
    cy = y + Mm(20)

    # 左側の大きな雲
    add_filled_oval(slide, cx - Mm(10), cy - Mm(7), Mm(15), Mm(12), RED)
    # 右側の小さな雲（左と少し重ねる）
    add_filled_oval(slide, cx + Mm(2),  cy - Mm(4), Mm(11), Mm(9),  RED)
    # 中央の白い穴
    add_filled_oval(slide, cx - Mm(1),  cy - Mm(2), Mm(5),  Mm(4),  WHITE)

    # 下部の "mamawell" テキスト
    add_text(
        slide, x + Mm(2), y + Mm(31), Mm(33), Mm(8),
        "mamawell", Pt(15), RED, bold=True,
        align=PP_ALIGN.CENTER, font_name="Arial Rounded MT Bold"
    )


def value_row(slide, x, y, top_offset, icon_glyph, jp, en, icon_bg):
    """3つの価値観の1行（アイコン円 + 日本語タイトル + 英語サブタイトル）"""
    icon_size = Mm(7)
    icon_x = x + Mm(40)
    icon_y = y + top_offset

    # アイコン円背景
    add_filled_oval(slide, icon_x, icon_y, icon_size, icon_size, icon_bg)

    # アイコングリフ（簡易）
    add_text(
        slide, icon_x, icon_y, icon_size, icon_size,
        icon_glyph, Pt(11), RED, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE
    )

    # テキスト
    text_x = icon_x + icon_size + Mm(2)
    text_w = x + CARD_W - Mm(2) - text_x

    add_text(
        slide, text_x, icon_y - Mm(0.5), text_w, Mm(4.5),
        jp, Pt(8.5), DARK, bold=True, font_name="游ゴシック"
    )
    add_text(
        slide, text_x, icon_y + Mm(3.8), text_w, Mm(3.5),
        en, Pt(6.5), RED, font_name="Helvetica Neue"
    )


def dotted_separator(slide, x, y, top_offset):
    """点線区切り"""
    line = slide.shapes.add_connector(
        1,
        x + Mm(40), y + top_offset,
        x + CARD_W - Mm(2), y + top_offset
    )
    line.line.color.rgb = LGRAY
    line.line.width = Pt(0.5)
    ln = line.line._get_or_add_ln()
    prstDash = etree.SubElement(ln, qn('a:prstDash'))
    prstDash.set('val', 'sysDot')


def draw_card(slide, x, y):
    card_background(slide, x, y)
    decorations(slide, x, y)
    logo_area(slide, x, y)

    # 3つの価値観
    value_row(slide, x, y, Mm(8),  "♡", "常に誠実に",       "Integrity First",       PINK)
    dotted_separator(slide, x, y, Mm(18))
    value_row(slide, x, y, Mm(21), "📈", "仕組みで勝つ",     "System > Hero",         BLUE_GRAY)
    dotted_separator(slide, x, y, Mm(31))
    value_row(slide, x, y, Mm(34), "👥", "前進のために話す", "Move Forward Together", PINK)


def trim_marks(slide):
    """カット位置を示すトンボ（四隅に小さな十字）"""
    mark = Mm(2)
    for r in range(ROWS):
        for c in range(COLS):
            left = MARGIN_X + c * (CARD_W + GAP_X)
            top  = MARGIN_Y + r * (CARD_H + GAP_Y)
            for cx, cy in [
                (left, top),
                (left + CARD_W, top),
                (left, top + CARD_H),
                (left + CARD_W, top + CARD_H),
            ]:
                h = slide.shapes.add_connector(1, cx - mark, cy, cx + mark, cy)
                h.line.color.rgb = LGRAY; h.line.width = Pt(0.25)
                v = slide.shapes.add_connector(1, cx, cy - mark, cx, cy + mark)
                v.line.color.rgb = LGRAY; v.line.width = Pt(0.25)


def main():
    prs = Presentation()
    prs.slide_width  = PAGE_W
    prs.slide_height = PAGE_H

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE

    for r in range(ROWS):
        for c in range(COLS):
            x = MARGIN_X + c * (CARD_W + GAP_X)
            y = MARGIN_Y + r * (CARD_H + GAP_Y)
            draw_card(slide, x, y)

    trim_marks(slide)

    out = '/home/user/hello-claude/mamawell_employee_card.pptx'
    prs.save(out)
    print(f"✓ 作成完了: {out}")
    if os.path.exists(LOGO_PATH):
        print(f"  ロゴ画像を埋め込みました: {LOGO_PATH}")
    else:
        print("  ロゴ画像なし → シェイプで代替表現を配置")
        print(f"  画像を {LOGO_PATH} に置いて再実行するとロゴを埋め込みます")


if __name__ == '__main__':
    main()
